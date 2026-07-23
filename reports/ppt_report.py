from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from core.metrics import kpis, monthly_summary, demand_summary, stop_frequency
from ai.conclusions import generate_conclusions
from ai.recommendations import generate_recommendations

def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(.5), Inches(.25), Inches(12.3), Inches(.65))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(18, 61, 118)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(.5), Inches(.9), Inches(12.3), Inches(.4))
        sub.text_frame.paragraphs[0].text = subtitle

def build_ppt(df) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Centro de Inteligencia Operacional KOA",
              f"Periodo {df['fecha'].min():%d/%m/%Y} - {df['fecha'].max():%d/%m/%Y}")
    text = slide.shapes.add_textbox(Inches(.8), Inches(2), Inches(11.7), Inches(3))
    p = text.text_frame.paragraphs[0]
    p.text = "Informe ejecutivo automático"
    p.font.size = Pt(30)
    p.font.bold = True

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Indicadores principales")
    metrics = kpis(df)
    labels = [
        ("Registros", metrics["registros"]), ("Días", metrics["dias"]),
        ("Usuarios", metrics["usuarios"]), ("Puntualidad", f"{metrics['puntual_0_5']:.1%}"),
        ("Efectivos tarde", metrics["efectivos_tarde"]), ("P90", f"{metrics['tiempo_efectivo_p90']:.1f} min"),
    ]
    for i, (label, value) in enumerate(labels):
        x = .7 + (i % 3) * 4.15
        y = 1.4 + (i // 3) * 2.3
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.6), Inches(1.5))
        tf = shape.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(16)
        p2 = tf.add_paragraph()
        p2.text = str(value)
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(18, 61, 118)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Hallazgos ejecutivos")
    box = slide.shapes.add_textbox(Inches(.7), Inches(1.2), Inches(12), Inches(5.6))
    tf = box.text_frame
    for i, item in enumerate(generate_conclusions(df)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(17)

    monthly = monthly_summary(df)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Demanda mensual")
    data = ChartData()
    data.categories = monthly["mes"].tolist()
    data.add_series("Usuarios", monthly["usuarios"].tolist())
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(.8), Inches(1.3), Inches(11.8), Inches(5.5), data)

    demand = demand_summary(df)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Demanda por recorrido - tarde")
    data = ChartData()
    data.categories = demand["recorrido"].tolist()
    data.add_series("Usuarios", demand["usuarios"].tolist())
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(.8), Inches(1.3), Inches(11.8), Inches(5.5), data)

    stops = stop_frequency(df)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Frecuencia efectiva de paraderos")
    data = ChartData()
    data.categories = stops["paradero"].tolist()
    data.add_series("Frecuencia", stops["frecuencia_efectiva"].fillna(0).tolist())
    slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(.8), Inches(1.3), Inches(11.8), Inches(5.5), data)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Plan de acción")
    box = slide.shapes.add_textbox(Inches(.6), Inches(1.1), Inches(12.2), Inches(5.9))
    tf = box.text_frame
    for i, rec in enumerate(generate_recommendations(df)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{rec['prioridad']} - {rec['accion']} ({rec['plazo']})"
        p.font.size = Pt(16)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
